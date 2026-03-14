"""Generate AI Efficiency Proposal slides as PowerPoint (.pptx).

Arctic-frost theme matching the HTML version.
Output: ai-efficiency-proposal.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE

# === THEME COLORS ===
FROST_900 = RGBColor(0x0C, 0x2D, 0x4A)
FROST_800 = RGBColor(0x13, 0x40, 0x68)
FROST_700 = RGBColor(0x1A, 0x52, 0x76)
FROST_600 = RGBColor(0x24, 0x71, 0xA3)
FROST_500 = RGBColor(0x2E, 0x86, 0xC1)
FROST_400 = RGBColor(0x5D, 0xAD, 0xE2)
FROST_300 = RGBColor(0x85, 0xC1, 0xE9)
FROST_200 = RGBColor(0xAE, 0xD6, 0xF1)
FROST_100 = RGBColor(0xD6, 0xEA, 0xF8)
FROST_50 = RGBColor(0xEA, 0xF2, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x1B, 0x26, 0x31)
TEXT_BODY = RGBColor(0x2C, 0x3E, 0x50)
TEXT_MUTED = RGBColor(0x5D, 0x6D, 0x7E)
TEAL = RGBColor(0x17, 0xA2, 0xB8)
EMERALD = RGBColor(0x1A, 0xBC, 0x9C)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
WARNING = RGBColor(0xF3, 0x9C, 0x12)
CORAL = RGBColor(0xE7, 0x4C, 0x3C)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide, color1, color2):
    """Add a gradient background using a full-slide shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(0), Emu(0),
        SLIDE_WIDTH, SLIDE_HEIGHT
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    shape.rotation = 0


def add_topbar(slide):
    """Add the colored top bar decoration."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        SLIDE_WIDTH, Pt(5)
    )
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = FROST_500
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = TEAL
    fill.gradient_stops[1].position = 1.0


def add_watermark(slide, text="Confidential"):
    txbox = slide.shapes.add_textbox(
        Inches(10.5), Inches(6.9), Inches(2.5), Inches(0.4)
    )
    tf = txbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_text(slide, left, top, width, height, text, font_size=14,
             color=TEXT_BODY, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Yu Gothic"):
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txbox


def add_multiline(slide, left, top, width, height, lines, font_size=13,
                  color=TEXT_BODY, line_spacing=1.5):
    """Add textbox with multiple paragraphs."""
    txbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Yu Gothic"
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txbox


def add_rounded_rect(slide, left, top, width, height, fill_color,
                     border_color=None, text="", font_size=12,
                     font_color=TEXT_BODY, bold=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.name = "Yu Gothic"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_metric_card(slide, left, top, width, value, label, sublabel="",
                    bg_color=FROST_50, border_color=FROST_200):
    """Add a metric card with large value + label."""
    card = add_rounded_rect(slide, left, top, width, 1.5,
                            bg_color, border_color)
    add_text(slide, left + 0.1, top + 0.2, width - 0.2, 0.7,
             value, font_size=32, color=FROST_600, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, left + 0.1, top + 0.85, width - 0.2, 0.3,
             label, font_size=11, color=TEXT_MUTED,
             alignment=PP_ALIGN.CENTER)
    if sublabel:
        add_text(slide, left + 0.1, top + 1.15, width - 0.2, 0.25,
                 sublabel, font_size=9, color=SUCCESS, bold=True,
                 alignment=PP_ALIGN.CENTER)
    return card


def add_bar(slide, left, top, width, pct, label, value_text,
            bar_color=FROST_500):
    """Add a horizontal bar chart row."""
    # Label
    add_text(slide, left, top, 1.5, 0.35, label,
             font_size=11, color=TEXT_BODY)
    # Track
    track_left = left + 1.6
    track_width = width - 2.8
    add_rounded_rect(slide, track_left, top + 0.05, track_width, 0.28,
                     FROST_100)
    # Fill
    fill_width = track_width * (pct / 100)
    if fill_width > 0.1:
        bar = add_rounded_rect(slide, track_left, top + 0.05,
                               fill_width, 0.28, bar_color)
        # Percent text inside bar
        add_text(slide, track_left, top + 0.05, fill_width - 0.05, 0.28,
                 f"{pct}%", font_size=9, color=WHITE, bold=True,
                 alignment=PP_ALIGN.RIGHT)
    # Value
    add_text(slide, left + width - 1.0, top, 0.9, 0.35,
             value_text, font_size=11, color=FROST_700, bold=True,
             alignment=PP_ALIGN.RIGHT)


def add_table(slide, left, top, width, height, data, col_widths=None):
    """Add a styled table. data = list of rows, first row is header."""
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            cell = table.cell(r, c)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Yu Gothic"

            if r == 0:
                # Header
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = FROST_800
            else:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_BODY
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = FROST_50
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE

    return table_shape


def add_insight_box(slide, left, top, width, text, accent_color=FROST_500):
    """Left-bordered highlight box."""
    # Accent bar
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top),
        Pt(4), Inches(0.8)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = accent_color
    slide.shapes[-1].line.fill.background()

    # Background
    add_rounded_rect(slide, left + 0.05, top, width - 0.05, 0.8,
                     FROST_50, FROST_200)
    add_text(slide, left + 0.2, top + 0.1, width - 0.4, 0.6,
             text, font_size=11, color=TEXT_BODY)


def add_step(slide, left, top, num, title, desc):
    """Add numbered step row."""
    # Circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top),
        Inches(0.35), Inches(0.35)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = FROST_500
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(num)
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = "Yu Gothic"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title + desc
    add_text(slide, left + 0.5, top - 0.02, 4.5, 0.22,
             title, font_size=13, color=TEXT_DARK, bold=True)
    add_text(slide, left + 0.5, top + 0.22, 4.5, 0.2,
             desc, font_size=10, color=TEXT_MUTED)


# ============================================================
# BUILD PRESENTATION
# ============================================================

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # blank

    # ── SLIDE 1: COVER ──
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, FROST_900, FROST_500)
    add_text(slide, 3, 1.2, 7.3, 0.5, "PROPOSAL 2026",
             font_size=14, color=FROST_300, alignment=PP_ALIGN.CENTER)
    add_text(slide, 2, 2.2, 9.3, 1.2, "AI活用による\n業務効率化提案",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Divider line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.8), Inches(3.8), Inches(1.7), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = FROST_300
    line.line.fill.background()
    add_text(slide, 3, 4.1, 7.3, 0.5, "経営企画部 御中",
             font_size=20, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text(slide, 3, 5.5, 7.3, 0.4, "2026年3月14日 | Confidential",
             font_size=12, color=FROST_300, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 2: EXECUTIVE SUMMARY ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "エグゼクティブサマリー",
             font_size=28, color=FROST_800, bold=True)

    add_insight_box(slide, 0.8, 1.3, 11.7,
                    "AI導入により年間約4,800時間の業務削減と1.2億円のコスト効果を見込む。"
                    "定型業務の自動化・データ分析の高度化・社内ナレッジ活用を3本柱に、6ヶ月で段階導入。")

    add_metric_card(slide, 0.8, 2.6, 3.6, "40%", "定型業務の削減率",
                    "前年比+18pt")
    add_metric_card(slide, 4.8, 2.6, 3.6, "1.2億", "年間コスト効果（円）",
                    "ROI 340%")
    add_metric_card(slide, 8.8, 2.6, 3.6, "6ヶ月", "導入完了まで",
                    "回収8ヶ月")
    add_watermark(slide)

    # ── SLIDE 3: SECTION — 現状分析 ──
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, FROST_800, FROST_600)
    add_text(slide, 9, 0.6, 4, 1.2, "01",
             font_size=72, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
             alignment=PP_ALIGN.RIGHT)
    # Make the 01 semi-transparent by using a lighter shade
    slide.shapes[-1].fill.background() if hasattr(slide.shapes[-1], 'fill') else None
    add_text(slide, 2, 2.8, 9.3, 1, "現状分析",
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, 2, 3.9, 9.3, 0.5, "業務課題の定量的把握",
             font_size=18, color=FROST_300, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 4: 部門別非効率 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "部門別 非効率業務の発生状況",
             font_size=28, color=FROST_800, bold=True)

    bars = [
        ("経理・財務", 78, "月62h"),
        ("営業・マーケ", 65, "月48h"),
        ("人事・総務", 58, "月40h"),
        ("経営企画", 52, "月35h"),
        ("情報システム", 45, "月28h"),
    ]
    for i, (label, pct, val) in enumerate(bars):
        add_bar(slide, 0.8, 1.5 + i * 0.55, 6.5, pct, label, val)

    # Donut representation (simplified as a circle + text)
    add_rounded_rect(slide, 8.2, 1.5, 4.2, 2.2, FROST_50, FROST_200)
    add_text(slide, 8.5, 1.6, 3.6, 0.4, "月間合計",
             font_size=11, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
    add_text(slide, 8.5, 2.0, 3.6, 0.6, "213時間",
             font_size=28, color=FROST_600, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline(slide, 8.5, 2.7, 3.6, 0.9, [
        "● 定型処理 54%",
        "● データ集計 21%",
        "● 情報検索 25%",
    ], font_size=10, color=TEXT_BODY)

    add_insight_box(slide, 0.8, 4.5, 11.7,
                    "全部門合計で月間213時間がAI自動化の対象。"
                    "年間換算で2,556時間の削減ポテンシャル。")
    add_watermark(slide)

    # ── SLIDE 5: 国内AI導入動向 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "国内企業のAI導入動向",
             font_size=28, color=FROST_800, bold=True)

    add_table(slide, 0.8, 1.3, 5.8, 1.6, [
        ["年度", "大企業", "中堅企業", "前年比"],
        ["2024", "42%", "18%", "—"],
        ["2025", "61%", "34%", "+19pt"],
        ["2026（予測）", "78%", "52%", "+17pt"],
    ], col_widths=[1.6, 1.2, 1.2, 1.0])

    add_text(slide, 7.2, 1.3, 5.2, 0.4, "効果実感の高い領域 TOP5",
             font_size=14, color=FROST_700, bold=True)
    top5 = [
        ("文書処理", 85), ("データ分析", 76), ("顧客対応", 68),
        ("ナレッジ共有", 62), ("コード生成", 54),
    ]
    for i, (label, pct) in enumerate(top5):
        add_bar(slide, 7.2, 1.9 + i * 0.48, 5.3, pct, label, f"{pct}%")

    add_insight_box(slide, 0.8, 4.2, 5.8,
                    "未導入リスク：競合との生産性格差が拡大。"
                    "2026年は「導入の最終ウィンドウ」。")

    add_rounded_rect(slide, 7.2, 4.2, 5.3, 0.8, RGBColor(0xE8, 0xF8, 0xF5),
                     TEAL)
    add_text(slide, 7.4, 4.3, 4.9, 0.6,
             "当社の注力3領域（文書・分析・ナレッジ）は\nいずれもTOP5に入る高効果領域",
             font_size=11, color=TEXT_BODY)
    add_watermark(slide, "出典：総務省 AI利活用実態調査 2025")

    # ── SLIDE 6: SECTION — AI活用戦略 ──
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, FROST_800, FROST_600)
    add_text(slide, 9, 0.6, 4, 1.2, "02",
             font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(slide, 2, 2.8, 9.3, 1, "AI活用戦略",
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, 2, 3.9, 9.3, 0.5, "3つの重点施策",
             font_size=18, color=FROST_300, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 7: 3本柱 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "業務効率化の3本柱",
             font_size=28, color=FROST_800, bold=True)

    pillars = [
        ("施策 1", "ドキュメント処理\n自動化",
         "請求書・契約書・レポートの\n読取・生成・要約をAIで自動化",
         "-40%", "処理時間削減", FROST_50, FROST_200),
        ("施策 2", "データ分析\n意思決定支援",
         "売上・KPIの自動集計と\nAIによるインサイト抽出",
         "3x", "分析速度向上",
         RGBColor(0xE8, 0xF8, 0xF5), RGBColor(0xA3, 0xE4, 0xD7)),
        ("施策 3", "社内ナレッジ\nAI活用",
         "社内文書・議事録・マニュアルを\nAIが横断検索・回答",
         "-65%", "情報検索時間",
         RGBColor(0xFE, 0xF9, 0xE7), RGBColor(0xF9, 0xE7, 0x9F)),
    ]

    for i, (tag, title, desc, val, metric, bg, border) in enumerate(pillars):
        x = 0.8 + i * 4.1
        add_rounded_rect(slide, x, 1.3, 3.7, 4.8, bg, border)
        # Tag
        add_rounded_rect(slide, x + 1.2, 1.5, 1.2, 0.32, FROST_100,
                         text=tag, font_size=9, font_color=FROST_700,
                         bold=True)
        # Title
        add_text(slide, x + 0.2, 2.1, 3.3, 0.8, title,
                 font_size=18, color=FROST_800, bold=True,
                 alignment=PP_ALIGN.CENTER)
        # Description
        add_text(slide, x + 0.2, 3.0, 3.3, 0.8, desc,
                 font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)
        # Metric
        add_text(slide, x + 0.2, 4.0, 3.3, 0.7, val,
                 font_size=30, color=FROST_600, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 4.7, 3.3, 0.3, metric,
                 font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

    add_watermark(slide)

    # ── SLIDE 8: 施策1 詳細 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "施策1：ドキュメント処理の自動化",
             font_size=28, color=FROST_800, bold=True)

    add_text(slide, 0.8, 1.2, 5.5, 0.4, "Before → After 比較",
             font_size=14, color=FROST_700, bold=True)
    add_table(slide, 0.8, 1.7, 5.8, 2.0, [
        ["業務", "現状", "AI導入後", "削減率"],
        ["請求書処理", "15分/件", "3分/件", "-80%"],
        ["契約書レビュー", "2時間/件", "30分/件", "-75%"],
        ["月次レポート", "8時間", "2時間", "-75%"],
        ["議事録作成", "45分/回", "5分/回", "-89%"],
    ], col_widths=[1.6, 1.2, 1.2, 1.0])

    add_insight_box(slide, 0.8, 4.0, 5.8,
                    "年間削減：約2,100時間 — FTE換算で1.2名分の人的リソースを創出")

    add_text(slide, 7.2, 1.2, 5.3, 0.4, "導入ツール構成",
             font_size=14, color=FROST_700, bold=True)
    add_rounded_rect(slide, 7.2, 1.7, 5.3, 1.6, FROST_50, FROST_200)
    add_text(slide, 7.5, 1.8, 4.7, 0.35, "Claude API + RAG基盤",
             font_size=13, color=FROST_600, bold=True)
    add_multiline(slide, 7.5, 2.2, 4.7, 1.0, [
        "・社内文書の読取・分類・要約",
        "・テンプレートベースの文書生成",
        "・複数文書の横断比較分析",
    ], font_size=10, color=TEXT_BODY)

    add_rounded_rect(slide, 7.2, 3.5, 5.3, 1.6, FROST_50, FROST_200)
    add_text(slide, 7.5, 3.6, 4.7, 0.35, "OCR + ワークフロー連携",
             font_size=13, color=FROST_600, bold=True)
    add_multiline(slide, 7.5, 4.0, 4.7, 1.0, [
        "・紙書類のデジタル化自動処理",
        "・既存システムとのAPI連携",
        "・承認フローの自動ルーティング",
    ], font_size=10, color=TEXT_BODY)
    add_watermark(slide)

    # ── SLIDE 9: 施策2 詳細 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "施策2：データ分析・意思決定支援",
             font_size=28, color=FROST_800, bold=True)

    add_text(slide, 0.8, 1.2, 5.5, 0.4, "現状の課題",
             font_size=14, color=FROST_700, bold=True)

    issues = [
        ("高", CORAL, "レポート作成に過大な工数", "月次レポートの集計・可視化に毎月40h以上"),
        ("高", CORAL, "データのサイロ化", "部門間でデータ定義・形式が不統一"),
        ("中", WARNING, "分析の属人化", "特定メンバーに依存した分析業務"),
    ]
    for i, (level, color, title, desc) in enumerate(issues):
        y = 1.7 + i * 0.8
        add_rounded_rect(slide, 0.8, y, 0.5, 0.3, color,
                         text=level, font_size=9, font_color=WHITE, bold=True)
        add_text(slide, 1.5, y - 0.02, 4.8, 0.25, title,
                 font_size=12, color=TEXT_DARK, bold=True)
        add_text(slide, 1.5, y + 0.25, 4.8, 0.25, desc,
                 font_size=10, color=TEXT_MUTED)

    add_text(slide, 7.2, 1.2, 5.3, 0.4, "AIソリューション",
             font_size=14, color=FROST_700, bold=True)

    steps = [
        ("1", "自動ダッシュボード", "KPI自動集計 → リアルタイム可視化"),
        ("2", "異常値の自動検知", "閾値超過や急変をSlack/メールで即通知"),
        ("3", "自然言語クエリ", "「先月の売上トップ3は？」で即回答"),
        ("4", "予測・シナリオ分析", "AIが複数シナリオを自動シミュレーション"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        add_step(slide, 7.2, 1.7 + i * 0.7, num, title, desc)

    add_rounded_rect(slide, 7.2, 4.7, 5.3, 0.5,
                     RGBColor(0xE8, 0xF8, 0xF5), TEAL)
    add_text(slide, 7.4, 4.8, 4.9, 0.35,
             "分析リードタイム：5日 → 当日（即時）",
             font_size=11, color=TEXT_BODY, bold=True)
    add_watermark(slide)

    # ── SLIDE 10: 施策3 詳細 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "施策3：社内ナレッジのAI活用",
             font_size=28, color=FROST_800, bold=True)

    # Chat mockup
    add_rounded_rect(slide, 0.8, 1.3, 5.8, 4.5, FROST_50, FROST_200)
    add_text(slide, 1.1, 1.4, 4, 0.3, "● 社内AIアシスタント — Online",
             font_size=10, color=TEXT_MUTED)

    # User bubble
    add_rounded_rect(slide, 1.2, 1.9, 4.2, 0.45,
                     RGBColor(0xE4, 0xE4, 0xE4),
                     text="出張旅費の精算ルールを教えて",
                     font_size=10, font_color=TEXT_DARK)
    # AI bubble
    add_rounded_rect(slide, 1.5, 2.5, 4.5, 0.85, FROST_500)
    add_text(slide, 1.7, 2.55, 4.1, 0.7,
             "出張旅費規程（2026年改定版）に基づき、\n"
             "日当は役職別に設定されています。\n"
             "部長以上：12,000円、課長：10,000円…",
             font_size=10, color=WHITE)
    # User bubble 2
    add_rounded_rect(slide, 1.2, 3.6, 4.8, 0.45,
                     RGBColor(0xE4, 0xE4, 0xE4),
                     text="前回の取締役会で決まった優先投資案件は？",
                     font_size=10, font_color=TEXT_DARK)
    # AI bubble 2
    add_rounded_rect(slide, 1.5, 4.3, 4.5, 0.55, FROST_500)
    add_text(slide, 1.7, 4.35, 4.1, 0.45,
             "3月度取締役会議事録より、\n以下の3案件が承認されました…",
             font_size=10, color=WHITE)

    # Right column
    add_text(slide, 7.2, 1.2, 5.3, 0.4, "構築アプローチ",
             font_size=14, color=FROST_700, bold=True)
    add_rounded_rect(slide, 7.2, 1.7, 5.3, 1.1, FROST_50, FROST_200)
    add_text(slide, 7.5, 1.8, 4.7, 0.3, "RAG（検索拡張生成）基盤",
             font_size=12, color=FROST_600, bold=True)
    add_text(slide, 7.5, 2.15, 4.7, 0.55,
             "社内文書をベクトルDBに格納。質問に対し\n関連情報を検索→AIが根拠付きで回答を生成。",
             font_size=10, color=TEXT_MUTED)

    add_rounded_rect(slide, 7.2, 3.0, 5.3, 1.2, FROST_50, FROST_200)
    add_text(slide, 7.5, 3.1, 4.7, 0.3, "対象データソース",
             font_size=12, color=FROST_600, bold=True)
    add_multiline(slide, 7.5, 3.45, 4.7, 0.7, [
        "・社内規程・マニュアル — 350件",
        "・議事録・報告書 — 年間2,000件",
        "・FAQ・問い合わせ履歴",
    ], font_size=10, color=TEXT_BODY)

    add_metric_card(slide, 7.2, 4.5, 2.5, "-65%", "検索時間削減",
                    bg_color=RGBColor(0xE8, 0xF8, 0xF5),
                    border_color=RGBColor(0xA3, 0xE4, 0xD7))
    add_metric_card(slide, 10.0, 4.5, 2.5, "-30%", "オンボーディング",
                    bg_color=RGBColor(0xE8, 0xF8, 0xF5),
                    border_color=RGBColor(0xA3, 0xE4, 0xD7))
    add_watermark(slide)

    # ── SLIDE 11: SECTION — 投資対効果 ──
    slide = prs.slides.add_slide(blank_layout)
    add_gradient_bg(slide, FROST_800, FROST_600)
    add_text(slide, 9, 0.6, 4, 1.2, "03",
             font_size=72, color=WHITE, bold=True, alignment=PP_ALIGN.RIGHT)
    add_text(slide, 2, 2.8, 9.3, 1, "投資対効果",
             font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, 2, 3.9, 9.3, 0.5, "ROI試算とロードマップ",
             font_size=18, color=FROST_300, alignment=PP_ALIGN.CENTER)

    # ── SLIDE 12: ROI試算 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "ROI試算（3年間）",
             font_size=28, color=FROST_800, bold=True)

    # ROI Flow
    flow_items = [
        ("初期投資", "3,200万", "開発・導入・研修", FROST_50, FROST_200, FROST_700),
        ("年間運用費", "960万", "API・保守・改善", FROST_50, FROST_200, FROST_700),
        ("年間効果", "1.2億", "工数削減+生産性向上", FROST_600, FROST_600, WHITE),
    ]
    for i, (label, val, sub, bg, border, text_c) in enumerate(flow_items):
        x = 0.8 + i * 4.3
        add_rounded_rect(slide, x, 1.3, 3.5, 1.5, bg, border)
        add_text(slide, x + 0.2, 1.35, 3.1, 0.3, label,
                 font_size=10, color=TEXT_MUTED if i < 2 else RGBColor(0xCC, 0xDD, 0xEE),
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 1.7, 3.1, 0.6, val,
                 font_size=26, color=text_c, bold=True,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 2.3, 3.1, 0.3, sub,
                 font_size=9, color=TEXT_MUTED if i < 2 else RGBColor(0xCC, 0xDD, 0xEE),
                 alignment=PP_ALIGN.CENTER)

    # Arrows
    for i in range(2):
        x = 4.5 + i * 4.3
        add_text(slide, x, 1.7, 0.6, 0.5, "▶",
                 font_size=18, color=FROST_300, alignment=PP_ALIGN.CENTER)

    # Table
    add_table(slide, 0.8, 3.2, 7.0, 2.0, [
        ["項目", "1年目", "2年目", "3年目"],
        ["コスト効果", "5,800万", "9,400万", "1.2億"],
        ["投資額", "4,160万", "960万", "960万"],
        ["純効果", "1,640万", "8,440万", "1.1億"],
        ["累計純効果", "1,640万", "1.01億", "2.11億"],
    ], col_widths=[2.0, 1.5, 1.5, 1.5])

    add_metric_card(slide, 8.3, 3.2, 2.2, "8ヶ月", "投資回収")
    add_metric_card(slide, 10.8, 3.2, 2.2, "340%", "3年累計ROI")
    add_watermark(slide)

    # ── SLIDE 13: ロードマップ ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "導入ロードマップ",
             font_size=28, color=FROST_800, bold=True)

    # Phase bar
    phases = [
        ("Phase 1", "1〜2ヶ月", True),
        ("Phase 2", "3〜4ヶ月", False),
        ("Phase 3", "5〜6ヶ月", False),
    ]
    for i, (label, period, active) in enumerate(phases):
        x = 0.8 + i * 4.1
        bg = FROST_500 if active else FROST_50
        tc = WHITE if active else TEXT_BODY
        add_rounded_rect(slide, x, 1.3, 3.7, 0.7, bg,
                         FROST_500 if active else FROST_200)
        add_text(slide, x + 0.2, 1.32, 3.3, 0.25, label,
                 font_size=9, color=RGBColor(0xBB,0xDD,0xEE) if active else TEXT_MUTED,
                 alignment=PP_ALIGN.CENTER)
        add_text(slide, x + 0.2, 1.55, 3.3, 0.35, period,
                 font_size=14, color=tc, bold=True, alignment=PP_ALIGN.CENTER)

    # Phase details
    phase_details = [
        ("基盤構築・PoC", [
            "AI基盤環境構築",
            "ドキュメント自動化PoC（経理）",
            "ナレッジDB設計・初期投入",
            "セキュリティ・ガバナンス整備",
        ], "PoC成功・効果実証", FROST_50, FROST_200),
        ("本格導入・展開", [
            "ドキュメント自動化を全部門展開",
            "データ分析ダッシュボード構築",
            "AIアシスタント β版リリース",
            "利用者トレーニング実施",
        ], "日常業務に定着",
         RGBColor(0xE8,0xF8,0xF5), RGBColor(0xA3,0xE4,0xD7)),
        ("最適化・発展", [
            "利用データに基づく精度改善",
            "予測分析・シナリオ機能追加",
            "外部データ連携の拡充",
            "次年度計画の策定",
        ], "ROI達成・拡大計画",
         RGBColor(0xFE,0xF9,0xE7), RGBColor(0xF9,0xE7,0x9F)),
    ]
    for i, (title, items, goal, bg, border) in enumerate(phase_details):
        x = 0.8 + i * 4.1
        add_rounded_rect(slide, x, 2.3, 3.7, 3.5, bg, border)
        add_text(slide, x + 0.2, 2.4, 3.3, 0.35, title,
                 font_size=13, color=FROST_700, bold=True)
        add_multiline(slide, x + 0.2, 2.8, 3.3, 2.0,
                      [f"・{item}" for item in items],
                      font_size=10, color=TEXT_BODY)
        add_text(slide, x + 0.2, 5.0, 3.3, 0.3,
                 f"ゴール：{goal}",
                 font_size=10, color=TEXT_MUTED, bold=True)
    add_watermark(slide)

    # ── SLIDE 14: リスクと対策 ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "想定リスクと対策",
             font_size=28, color=FROST_800, bold=True)

    add_table(slide, 0.8, 1.3, 11.7, 3.0, [
        ["重要度", "リスク", "影響", "対策", "担当"],
        ["高", "情報セキュリティ", "機密データの外部漏洩",
         "VPC環境、暗号化、アクセス制御の多層化", "情報システム部"],
        ["高", "AIの回答精度", "ハルシネーションによる誤情報",
         "RAG根拠明示、人間レビュー、信頼度スコア", "AI推進チーム"],
        ["中", "社内定着", "利用率低迷で効果限定的",
         "チェンジマネジメント、部門推進者配置", "経営企画部"],
        ["中", "コスト超過", "API利用量の想定超過",
         "利用量監視、キャッシュ戦略、月次レビュー", "経理・情シス"],
        ["低", "技術陳腐化", "モデル更新による設計変更",
         "マルチモデル抽象化、四半期技術レビュー", "AI推進チーム"],
    ], col_widths=[1.0, 1.8, 2.5, 3.5, 1.5])

    add_insight_box(slide, 0.8, 4.8, 11.7,
                    "リスク軽減方針：Phase 1のPoC期間でセキュリティ・精度検証を完了し、"
                    "本格展開前にガバナンス体制を確立。")
    add_watermark(slide)

    # ── SLIDE 15: まとめ ──
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, WHITE)
    add_topbar(slide)
    add_text(slide, 0.8, 0.5, 11, 0.6, "まとめ・ネクストステップ",
             font_size=28, color=FROST_800, bold=True)

    # Summary metrics
    add_text(slide, 0.8, 1.2, 5.5, 0.4, "提案サマリー",
             font_size=16, color=FROST_700, bold=True)
    metrics = [("4,800h", "年間削減"), ("1.2億", "年間効果"),
               ("340%", "3年ROI"), ("8ヶ月", "投資回収")]
    for i, (val, label) in enumerate(metrics):
        x = 0.8 + i * 1.5
        add_metric_card(slide, x, 1.7, 1.3, val, label)

    add_insight_box(slide, 0.8, 3.6, 5.8,
                    "AI導入は「コスト削減」と「意思決定の質・速度向上」の"
                    "両方を実現する戦略投資。早期着手が競争優位の源泉。")

    # Next steps
    add_text(slide, 7.2, 1.2, 5.3, 0.4, "ネクストステップ",
             font_size=16, color=FROST_700, bold=True)

    next_steps = [
        ("1", "経営会議での承認", "本提案の方針決定 — 3月中"),
        ("2", "プロジェクト体制構築", "推進チーム発足・部門推進者選定"),
        ("3", "PoC対象業務の確定", "経理部門のドキュメント処理を推奨"),
        ("4", "Phase 1 キックオフ", "4月第2週を目標"),
    ]
    for i, (num, title, desc) in enumerate(next_steps):
        add_step(slide, 7.2, 1.7 + i * 0.7, num, title, desc)

    # CTA
    add_rounded_rect(slide, 7.2, 4.7, 5.3, 0.8, FROST_800)
    add_text(slide, 7.5, 4.8, 4.7, 0.3,
             "ご質問・ご要望をお聞かせください",
             font_size=14, color=FROST_200, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(slide, 7.5, 5.15, 4.7, 0.25,
             "経営企画部 AI推進担当",
             font_size=9, color=FROST_400, alignment=PP_ALIGN.CENTER)
    add_watermark(slide)

    # Save
    output_path = "/home/user/kaizen-dx-lab/projects/airpro-duct/ai-efficiency-proposal.pptx"
    prs.save(output_path)
    print(f"Generated: {output_path}")
    return output_path


if __name__ == "__main__":
    build()
